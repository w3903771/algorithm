"""把三处本地资料（C++ 算法代码 / 信息学竞赛课件 / Pascal 模板）抽取成 Markdown 落到 sources/ 下。

用法: uv run python scripts/extract_local.py
"""
from __future__ import annotations

import os
import re
import sys
import zipfile
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
SOURCES = ROOT / "sources"

SRC_CPP = Path(r"E:\Programming\C++\code(c++)\Algorithm&DataStructure")
SRC_OI = Path(r"E:\Programming\C++\信息学竞赛课件（2018级）")
SRC_PAS = Path(r"E:\Programming\Pascal")

CODE_EXT = {".cpp": "cpp", ".c": "c", ".h": "cpp", ".pas": "pascal", ".py": "python"}
TEXT_EXT = {".txt", ".ini", ".dat"}

ENCODINGS = ("utf-8-sig", "gbk", "utf-16", "latin-1")


def read_text(path: Path) -> str:
    raw = path.read_bytes()
    for enc in ENCODINGS:
        try:
            return raw.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return raw.decode("utf-8", errors="replace")


def slug(name: str) -> str:
    return re.sub(r'[\\/:*?"<>|]+', "_", name).strip()


# ---------------------------------------------------------------- 提取器


def extract_pdf(path: Path) -> str:
    import fitz  # pymupdf

    out = []
    with fitz.open(path) as doc:
        for i, page in enumerate(doc, 1):
            text = page.get_text("text").strip()
            if text:
                out.append(f"### [第 {i} 页]\n\n{text}")
    return "\n\n".join(out)


def extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        chunks = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    chunks.append(t)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    chunks.append(" | ".join(c.text.strip() for c in row.cells))
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                chunks.append(f"[备注] {note}")
        if chunks:
            out.append(f"### [第 {i} 页]\n\n" + "\n\n".join(chunks))
    return "\n\n".join(out)


def extract_docx(path: Path) -> str:
    import docx

    d = docx.Document(str(path))
    out = []
    for p in d.paragraphs:
        t = p.text.rstrip()
        if t:
            style = (p.style.name or "").lower()
            if style.startswith("heading"):
                lvl = "".join(ch for ch in style if ch.isdigit()) or "2"
                out.append("#" * min(int(lvl) + 2, 6) + " " + t)
            else:
                out.append(t)
    for tbl in d.tables:
        for row in tbl.rows:
            out.append(" | ".join(c.text.strip().replace("\n", " ") for c in row.cells))
    return "\n\n".join(out)


def extract_ppt_or_doc_via_word(path: Path) -> str:
    """老格式 .doc/.ppt 用 COM 转成 txt。没有 Office 就返回空。"""
    try:
        import win32com.client  # type: ignore
    except ImportError:
        return ""
    tmp = ROOT / "sources" / "_tmp"
    tmp.mkdir(parents=True, exist_ok=True)
    out_file = tmp / (slug(path.stem) + ".txt")
    try:
        if path.suffix.lower() == ".doc":
            app = win32com.client.Dispatch("Word.Application")
            app.Visible = False
            doc = app.Documents.Open(str(path), ReadOnly=True)
            doc.SaveAs(str(out_file), FileFormat=2)  # wdFormatText
            doc.Close(False)
            app.Quit()
        else:
            return ""
    except Exception as exc:  # pragma: no cover
        print(f"    [warn] COM 转换失败 {path.name}: {exc}")
        return ""
    if out_file.exists():
        return read_text(out_file)
    return ""


def extract_zip_listing(path: Path) -> str:
    try:
        with zipfile.ZipFile(path) as z:
            return "\n".join(f"- {n}" for n in z.namelist())
    except Exception:
        return ""


# ---------------------------------------------------------------- 驱动


def handle(path: Path, base: Path, out_dir: Path, index: list) -> None:
    ext = path.suffix.lower()
    rel = path.relative_to(base)
    target = out_dir / rel.parent / (slug(path.stem) + ".md")

    body = ""
    kind = ""
    if ext in CODE_EXT:
        kind = "源码"
        body = f"```{CODE_EXT[ext]}\n{read_text(path).rstrip()}\n```"
    elif ext in TEXT_EXT:
        kind = "文本"
        body = f"```\n{read_text(path).rstrip()}\n```"
    elif ext == ".pdf":
        kind = "PDF 课件"
        body = extract_pdf(path)
    elif ext == ".pptx":
        kind = "PPT 课件"
        body = extract_pptx(path)
    elif ext == ".docx":
        kind = "Word 文档"
        body = extract_docx(path)
    elif ext == ".doc":
        kind = "Word 文档(旧)"
        body = extract_ppt_or_doc_via_word(path)
    elif ext in (".zip", ".rar"):
        kind = "压缩包"
        body = extract_zip_listing(path)
    else:
        return

    if not body.strip():
        print(f"    [skip] 无文本: {rel}")
        return

    target.parent.mkdir(parents=True, exist_ok=True)
    header = f"# {path.stem}\n\n> 来源: `{path}`  \n> 类型: {kind}\n\n---\n\n"
    target.write_text(header + body + "\n", encoding="utf-8")
    index.append((str(rel), kind, target.relative_to(SOURCES).as_posix(), len(body)))
    print(f"    [ok] {rel} -> {target.relative_to(ROOT)} ({len(body)} 字符)")


def walk(base: Path, out_name: str, only_files: list = None) -> list:
    out_dir = SOURCES / out_name
    out_dir.mkdir(parents=True, exist_ok=True)
    index: list = []
    print(f"\n== 处理 {base}")
    files = [base / f for f in only_files] if only_files else sorted(base.rglob("*"))
    for p in files:
        if p.is_file():
            try:
                handle(p, base, out_dir, index)
            except Exception as exc:
                print(f"    [err] {p.name}: {exc}")
    # 目录索引
    lines = ["# 提取索引\n", f"> 原始目录: `{base}`\n", "", "| 原文件 | 类型 | 提取结果 | 字符数 |", "| --- | --- | --- | --- |"]
    for rel, kind, tgt, n in index:
        lines.append(f"| {rel} | {kind} | [{Path(tgt).name}]({Path(tgt).relative_to(out_name).as_posix()}) | {n} |")
    (out_dir / "INDEX.md").write_text("\n".join(lines) + "\n", encoding="utf-8")
    return index


def main() -> int:
    SOURCES.mkdir(parents=True, exist_ok=True)
    total = []
    total += walk(SRC_CPP, "01-cpp-algo-ds")
    total += walk(SRC_OI, "02-oi-courseware")
    pas_files = [
        "模板.docx", "insertsort.pas", "jingzhuan.pas", "jishu.pas", "maopa.pas",
        "maopaoyouhua.pas", "noip2015跳石子.pas", "quicksort.pas", "sort.pas",
        "Sort函数.cpp", "xuanze.pas", "zhizhang.pas", "ziben.pas",
    ]
    total += walk(SRC_PAS, "03-pascal-template", only_files=pas_files)
    print(f"\n共提取 {len(total)} 个文件到 {SOURCES}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
