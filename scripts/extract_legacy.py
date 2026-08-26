"""处理老格式 .doc / .ppt（走 Office COM），输出到 sources/02-oi-courseware/day6/DP资料/。

用法: uv run python scripts/extract_legacy.py
"""
from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "sources" / "02-oi-courseware" / "day6" / "DP资料"
UNRAR = ROOT / "sources" / "_unrar"

TARGETS = [
    (UNRAR / "DP资料/DP资料/基础资料/背包九讲完整版.doc", OUT / "基础资料/背包九讲完整版.md"),
    (UNRAR / "DP的常见优化/DP的常见优化/动规问题的优化.pptx", OUT / "进阶资料/动规问题的优化.md"),
    (UNRAR / "DP的常见优化/DP的常见优化/source/单调队列.pptx", OUT / "进阶资料/单调队列.md"),
    (UNRAR / "DP的常见优化/DP的常见优化/source/浅析1D1D动态规划的优化(zzx).ppt", OUT / "进阶资料/浅析1D1D动态规划的优化.md"),
    (UNRAR / "DP的常见优化/DP的常见优化/source/8.杨哲《凸完全单调性的一个加强与应用》.ppt", OUT / "进阶资料/凸完全单调性的加强与应用.md"),
]


def doc_to_text(path: Path) -> str:
    import win32com.client as com

    app = com.gencache.EnsureDispatch("Word.Application")
    app.Visible = False
    tmp = ROOT / "sources" / "_tmp"
    tmp.mkdir(parents=True, exist_ok=True)
    out = tmp / (path.stem + ".txt")
    doc = app.Documents.Open(str(path), ReadOnly=True, ConfirmConversions=False)
    doc.SaveAs2(str(out), FileFormat=7)  # wdFormatEncodedText
    doc.Close(False)
    app.Quit()
    for enc in ("utf-8-sig", "utf-16", "gbk", "latin-1"):
        try:
            return out.read_bytes().decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return ""


def ppt_to_text(path: Path) -> str:
    """.ppt 用 PowerPoint COM 先另存为 .pptx，再用 python-pptx 解析。"""
    import win32com.client as com
    from pptx import Presentation

    tmp = ROOT / "sources" / "_tmp"
    tmp.mkdir(parents=True, exist_ok=True)
    newp = tmp / (path.stem + ".pptx")
    if not newp.exists():
        app = com.gencache.EnsureDispatch("PowerPoint.Application")
        pres = app.Presentations.Open(str(path), WithWindow=False)
        pres.SaveAs(str(newp), 24)  # ppSaveAsOpenXMLPresentation
        pres.Close()
        app.Quit()
    return pptx_text(newp)


def pptx_text(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        chunks = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                chunks.append(shape.text_frame.text.strip())
        if slide.has_notes_slide:
            n = slide.notes_slide.notes_text_frame.text.strip()
            if n:
                chunks.append(f"[备注] {n}")
        if chunks:
            out.append(f"### [第 {i} 页]\n\n" + "\n\n".join(chunks))
    return "\n\n".join(out)


def main() -> int:
    for src, dst in TARGETS:
        if not src.exists():
            print(f"[miss] {src}")
            continue
        ext = src.suffix.lower()
        try:
            if ext == ".doc":
                body = doc_to_text(src)
            elif ext == ".ppt":
                body = ppt_to_text(src)
            else:
                body = pptx_text(src)
        except Exception as exc:
            print(f"[err] {src.name}: {exc}")
            continue
        if not body.strip():
            print(f"[skip] {src.name} 无文本")
            continue
        dst.parent.mkdir(parents=True, exist_ok=True)
        dst.write_text(f"# {src.stem}\n\n> 来源: `{src.name}`（自 rar 解包）\n\n---\n\n{body}\n", encoding="utf-8")
        print(f"[ok] {src.name} -> {dst.relative_to(ROOT)} ({len(body)} 字符)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
